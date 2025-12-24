import os
import glob
import time
import sqlite3
import io
from typing import List

from sentence_transformers import SentenceTransformer
import weaviate
from weaviate.collections.classes.config import Property, DataType, Configure
from weaviate.collections.classes.filters import Filter

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

import pytesseract
from PIL import Image

# ---------------- CONFIG ----------------

CLASS_NAME = "Documents"
INDEX_DB = "index_state.db"

ALLOWED_EXT = (".txt", ".pdf", ".docx", ".ppt", ".pptx")

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
MIN_CHUNK_LEN = 40

# OCR Configuration
ENABLE_OCR = True
OCR_WORD_THRESHOLD = 50  # OCR triggers if extracted text has fewer than this many words
OCR_MAX_PAGES = 5        # Maximum pages to OCR per PDF

# ----------------------------------------


# -------- FILE READERS --------

def read_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except:
        return ""


def read_pdf(path: str) -> str:
    """
    Text-first PDF extraction with conditional OCR fallback.
    """
    try:
        doc = fitz.open(path)
    except:
        return ""

    texts = []

    # ---- PASS 1: TEXT EXTRACTION ----
    for page in doc:
        text = page.get_text().strip()
        if text:
            texts.append(text)

    full_text = "\n".join(texts)

    # ---- PASS 2: CONDITIONAL OCR ----
    word_count = len(full_text.split())

    if word_count >= OCR_WORD_THRESHOLD:
        return full_text

    # OCR only if text extraction failed AND OCR is enabled
    if not ENABLE_OCR:
        return full_text

    print(f"[OCR] Triggered for {os.path.basename(path)} (only {word_count} words extracted)")
    ocr_texts = []

    for page_num, page in enumerate(doc):
        # Stop after max pages limit
        if page_num >= OCR_MAX_PAGES:
            break

        try:
            pix = page.get_pixmap(dpi=300)
            # Use a more reliable temp file approach
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))
            ocr_text = pytesseract.image_to_string(img)
            if ocr_text.strip():
                ocr_texts.append(ocr_text)
        except Exception as e:
            print(f"[WARN] OCR failed for page {page_num + 1}: {e}")
            continue

        # Early termination: stop if we've extracted enough text
        current_ocr_text = "\n".join(ocr_texts)
        if len(current_ocr_text.split()) >= OCR_WORD_THRESHOLD:
            break

    ocr_result = "\n".join(ocr_texts)
    
    # IMPORTANT: Combine both text extraction and OCR results
    combined = full_text + "\n" + ocr_result if full_text.strip() else ocr_result
    
    final_word_count = len(combined.split())
    print(f"[OCR] Complete: {final_word_count} words total")
    
    return combined


def read_docx(path: str) -> str:
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except:
        return ""


def read_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        texts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            texts.append(p.text)

        return "\n".join(texts)
    except:
        return ""


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".txt":
        return read_txt(path)
    if ext == ".pdf":
        return read_pdf(path)
    if ext == ".docx":
        return read_docx(path)
    if ext in (".ppt", ".pptx"):
        return read_pptx(path)
    return ""


# -------- CHUNKING --------

def chunk_text(text: str) -> List[str]:
    text = " ".join(text.split())
    chunks = []
    start = 0

    while start < len(text):
        end = start + CHUNK_SIZE
        chunk = text[start:end].strip()
        if len(chunk) >= MIN_CHUNK_LEN:
            chunks.append(chunk)
        start = end - CHUNK_OVERLAP

    return chunks


# -------- SQLITE --------

def init_db():
    conn = sqlite3.connect(INDEX_DB)
    cur = conn.cursor()

    cur.execute("""
        CREATE TABLE IF NOT EXISTS indexed_files (
            path TEXT PRIMARY KEY,
            mtime REAL NOT NULL,
            size INTEGER NOT NULL,
            indexed_at REAL NOT NULL
        )
    """)

    cur.execute("""
        CREATE TABLE IF NOT EXISTS user_roots (
            path TEXT PRIMARY KEY
        )
    """)

    conn.commit()
    return conn


def load_user_roots(cur):
    cur.execute("SELECT path FROM user_roots")
    return [row[0] for row in cur.fetchall()]


# -------- MAIN INDEXER --------

def main():
    print("[INIT] Loading embedding model...")
    model = SentenceTransformer("all-MiniLM-L6-v2")
    print("[DONE] Model loaded")

    print("[INIT] Connecting to Weaviate...")
    client = weaviate.connect_to_local()
    print("[DONE] Connected")

    print("[INIT] Ensuring schema...")
    if CLASS_NAME not in client.collections.list_all():
        client.collections.create(
            name=CLASS_NAME,
            properties=[
                Property(name="file", data_type=DataType.TEXT),
                Property(name="path", data_type=DataType.TEXT),
                Property(name="chunk", data_type=DataType.TEXT),
            ],
            vector_config=Configure.Vectors.none(),
        )
        print("[DONE] Schema created")
    else:
        print("[INFO] Schema already exists")

    collection = client.collections.get(CLASS_NAME)

    conn = init_db()
    cur = conn.cursor()

    roots = load_user_roots(cur)
    if not roots:
        print("[WARN] No user roots configured. Indexer exiting.")
        conn.close()
        client.close()
        return

    print(f"[INFO] Using {len(roots)} user-defined roots")

    all_files = set()
    for root in roots:
        if not os.path.exists(root):
            continue
        for ext in ALLOWED_EXT:
            pattern = os.path.join(root, "**", f"*{ext}")
            for f in glob.glob(pattern, recursive=True):
                all_files.add(os.path.abspath(f))

    print(f"[INFO] Found {len(all_files)} files")

    cur.execute("SELECT path, mtime, size FROM indexed_files")
    known = {row[0]: (row[1], row[2]) for row in cur.fetchall()}

    to_index = []
    for path in all_files:
        try:
            stat = os.stat(path)
        except:
            continue

        if path not in known:
            to_index.append(path)
        else:
            old_mtime, old_size = known[path]
            if stat.st_mtime != old_mtime or stat.st_size != old_size:
                to_index.append(path)

    deleted = set(known.keys()) - all_files

    print(f"[INFO] New/changed files: {len(to_index)}")
    print(f"[INFO] Deleted files: {len(deleted)}")

    for path in deleted:
        collection.data.delete_many(
            where=Filter.by_property("path").equal(path)
        )
        cur.execute("DELETE FROM indexed_files WHERE path=?", (path,))

    for path in to_index:
        print(f"[INDEX] Processing: {os.path.basename(path)}")
        text = extract_text(path)
        if not text.strip():
            print(f"[WARN] No text extracted from {os.path.basename(path)}")
            continue

        chunks = chunk_text(text)
        if not chunks:
            print(f"[WARN] No chunks created from {os.path.basename(path)}")
            continue

        print(f"[INDEX] Created {len(chunks)} chunks ({len(text.split())} words)")
        vectors = model.encode(chunks)

        collection.data.delete_many(
            where=Filter.by_property("path").equal(path)
        )

        for chunk, vec in zip(chunks, vectors):
            collection.data.insert(
                properties={
                    "file": os.path.basename(path),
                    "path": path,
                    "chunk": chunk,
                },
                vector=vec.tolist(),
            )

        stat = os.stat(path)
        cur.execute(
            "REPLACE INTO indexed_files VALUES (?, ?, ?, ?)",
            (path, stat.st_mtime, stat.st_size, time.time())
        )

    conn.commit()
    conn.close()
    client.close()

    print("[DONE] Indexing complete")


if __name__ == "__main__":
    main()
